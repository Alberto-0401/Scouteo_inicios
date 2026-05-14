# -*- coding: utf-8 -*-
"""SCOUTEO TFG - Presentacion PowerPoint PRO"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Design System ────────────────────────────────────────────────────────────
BG    = RGBColor(0x07, 0x0C, 0x18)   # negro azulado profundo
BG2   = RGBColor(0x0E, 0x16, 0x2B)   # card oscura
BLUE  = RGBColor(0x25, 0x63, 0xEB)   # azul primario
LBLUE = RGBColor(0x60, 0xA5, 0xFA)   # azul claro
GOLD  = RGBColor(0xF5, 0x9E, 0x0B)   # dorado
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xD1, 0xD5, 0xDB)
GRAY  = RGBColor(0x6B, 0x72, 0x80)
DGRAY = RGBColor(0x37, 0x41, 0x51)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED   = RGBColor(0xEF, 0x44, 0x44)
PURP  = RGBColor(0x7C, 0x3A, 0xED)
TEAL  = RGBColor(0x06, 0xB6, 0xD4)
DGRN  = RGBColor(0x05, 0x7A, 0x55)
WCARD = RGBColor(0xF0, 0xF4, 0xFF)   # blanco calido para tarjetas de diagrama
F = "Segoe UI"

W = Inches(13.33)
H = Inches(7.5)
ARCH = r"C:\Users\User\Downloads\arquitectura.drawio.png"
ER   = r"C:\Users\User\Downloads\E_R.drawio.png"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── Utilidades ───────────────────────────────────────────────────────────────
def S():
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(sl, fill, l, t, w, h, lc=None, lw=Pt(1)):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else:   sh.line.fill.background()
    return sh

def O(sl, fill, l, t, w, h):            # ovalo / circulo
    sh = sl.shapes.add_shape(9, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def T(sl, text, l, t, w, h, sz=18, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.name = F; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tf

def Tp(tf, text, sz=15, bold=False, color=WHITE, before=8, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.space_before = Pt(before); p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = F
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return p

def IMG(sl, path, l, t, max_w, max_h=None):
    """Inserta imagen centrada y escalada para caber en max_w x max_h sin distorsion."""
    if not os.path.exists(path):
        return False
    from PIL import Image as _PIL
    with _PIL.open(path) as img:
        iw, ih = img.size
    aspect = iw / ih
    if max_h is None:
        sl.shapes.add_picture(path, l, t, max_w)
    else:
        w = max_w
        h = w / aspect
        if h > max_h:
            h = max_h
            w = h * aspect
        off_x = (max_w - w) / 2
        off_y = (max_h - h) / 2
        sl.shapes.add_picture(path, l + off_x, t + off_y, w, h)
    return True

def bg(sl):
    R(sl, BG, 0, 0, W, H)

def branding(sl):
    T(sl, "SCOUTEO", W - Inches(1.9), Inches(0.08), Inches(1.8), Inches(0.38),
      sz=10, bold=True, color=DGRAY, align=PP_ALIGN.RIGHT)

def slide_num(sl, n, total=7):
    T(sl, f"{n} / {total}", W - Inches(0.9), H - Inches(0.42),
      Inches(0.75), Inches(0.32), sz=10, color=DGRAY, align=PP_ALIGN.RIGHT)

def gold_line(sl, l, t, w=Inches(0.8), h=Inches(0.035)):
    R(sl, GOLD, l, t, w, h)

def header_block(sl, tag, title, sub=None):
    """Tag pill + gran titulo + linea dorada + subtitulo opcional"""
    # tag pill
    pw = Inches(len(tag) * 0.14 + 0.45)
    R(sl, BLUE, Inches(0.55), Inches(0.32), pw, Inches(0.34))
    T(sl, tag.upper(), Inches(0.55), Inches(0.31), pw, Inches(0.34),
      sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # title
    T(sl, title, Inches(0.55), Inches(0.72), W - Inches(1.1), Inches(0.95),
      sz=40, bold=True, color=WHITE)
    # gold underline
    gold_line(sl, Inches(0.55), Inches(1.6))
    if sub:
        T(sl, sub, Inches(0.55), Inches(1.7), W - Inches(2), Inches(0.45),
          sz=15, color=GRAY)

def num_circle(sl, num_str, cx, cy, r=Inches(0.42), fill=BLUE):
    """Circulo numerado centrado en (cx, cy)"""
    O(sl, fill, cx - r/2, cy - r/2, r, r)
    T(sl, num_str, cx - r/2, cy - r/2 - Inches(0.02), r, r,
      sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def add_transition(slide, kind="fade"):
    """Inyecta transicion OOXML en el slide. kind: fade|push|cover|wipe"""
    sld_el = slide._element
    for old in sld_el.findall(qn('p:transition')):
        sld_el.remove(old)
    trans = etree.SubElement(sld_el, qn('p:transition'))
    trans.set('spd', 'med')
    if kind == "fade":
        etree.SubElement(trans, qn('p:fade'))
    elif kind == "push":
        push = etree.SubElement(trans, qn('p:push'))
        push.set('dir', 'l')
    elif kind == "cover":
        cover = etree.SubElement(trans, qn('p:cover'))
        cover.set('dir', 'l')
    elif kind == "wipe":
        wipe = etree.SubElement(trans, qn('p:wipe'))
        wipe.set('dir', 'l')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)

# Franja izquierda dorada
R(sl, GOLD, 0, 0, Inches(0.06), H)

# Franja superior azul muy fina
R(sl, BLUE, 0, 0, W, Inches(0.04))

# Nombre del proyecto enorme
T(sl, "SCOUTEO",
  Inches(0.5), Inches(1.1), Inches(10), Inches(2.2),
  sz=96, bold=True, color=WHITE)

# Linea dorada bajo el titulo
gold_line(sl, Inches(0.55), Inches(3.2), Inches(1.4), Inches(0.05))

# Tagline
T(sl, "Sistema de Gestion Deportiva para Clubes de Futbol",
  Inches(0.55), Inches(3.35), Inches(9), Inches(0.65),
  sz=22, color=LGRAY)

# Separador
R(sl, DGRAY, Inches(0.55), Inches(4.15), Inches(7), Inches(0.02))

# Info autor
T(sl, "Alberto Perez Oncina",
  Inches(0.55), Inches(4.3), Inches(7), Inches(0.52),
  sz=20, bold=True, color=WHITE)
T(sl, "CFGS Desarrollo de Aplicaciones Multiplataforma  |  Curso 2024 - 2025",
  Inches(0.55), Inches(4.82), Inches(9.5), Inches(0.45),
  sz=14, color=GRAY)

# Chips de tech (esquina inferior derecha, apilados vertical)
tech = [
    ("JavaFX 25",    BLUE),
    ("Spring Boot",  DGRN),
    ("PostgreSQL",   TEAL),
    ("Cloud Run",    PURP),
    ("React Native", RGBColor(0x92, 0x40, 0x00)),
]
for i, (lbl, col) in enumerate(tech):
    lx = Inches(10.6)
    ty = Inches(2.2 + i * 0.88)
    R(sl, col, lx, ty, Inches(2.45), Inches(0.62))
    T(sl, lbl, lx, ty + Inches(0.08), Inches(2.45), Inches(0.5),
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_transition(sl, "fade")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 2)
header_block(sl, "Contexto", "Por que Scouteo?",
             "Los clubes de base gestionan de forma manual y dispersa")

cards = [
    (GOLD,  "01", "Comunicacion",
     "Convocatorias, avisos y cambios de horario\npor WhatsApp — sin historial ni trazabilidad."),
    (RED,   "02", "Datos",
     "Estadisticas de jugadores y plantillas\ngestionadas en Excel, sin conexion entre modulos."),
    (PURP,  "03", "Gestion medica",
     "Lesiones e historial medico en papel —\nimposible de consultar rapidamente."),
]

for i, (col, num, title, desc) in enumerate(cards):
    lx = Inches(0.55 + i * 4.2)
    ty = Inches(2.1)
    cw = Inches(3.9)
    ch = Inches(4.7)
    # Card de fondo
    R(sl, BG2, lx, ty, cw, ch, lc=DGRAY, lw=Pt(0.75))
    # Barra superior de color
    R(sl, col, lx, ty, cw, Inches(0.08))
    # Numero grande
    T(sl, num, lx + Inches(0.25), ty + Inches(0.2), cw - Inches(0.3), Inches(1.0),
      sz=52, bold=True, color=col)
    # Titulo
    T(sl, title, lx + Inches(0.25), ty + Inches(1.1), cw - Inches(0.3), Inches(0.55),
      sz=18, bold=True, color=WHITE)
    # Descripcion
    T(sl, desc, lx + Inches(0.25), ty + Inches(1.65), cw - Inches(0.3), Inches(2.0),
      sz=14, color=LGRAY)

# Solucion
R(sl, BLUE, Inches(0.55), Inches(7.0), Inches(12.33), Inches(0.38))
T(sl, "Scouteo centraliza jugadores, partidos, entrenamientos, lesiones e informes en una sola plataforma.",
  Inches(0.55), Inches(7.02), Inches(12.33), Inches(0.36),
  sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_transition(sl, "push")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITECTURA
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 3)
header_block(sl, "Diseno", "Arquitectura del Sistema")

# Diagrama en carta blanca (derecha)
R(sl, WCARD, Inches(5.8), Inches(1.8), Inches(7.1), Inches(5.35))
if not IMG(sl, ARCH, Inches(5.85), Inches(1.85), Inches(7.0), Inches(5.25)):
    T(sl, "[ arquitectura.drawio.png ]",
      Inches(5.8), Inches(4.2), Inches(7.1), Inches(0.5),
      sz=11, color=GRAY, align=PP_ALIGN.CENTER)

# Componentes (izquierda)
comps = [
    (GOLD,  "App Escritorio",  "JavaFX 25 + FXML + jpackage (.exe)"),
    (LBLUE, "API REST",        "Spring Boot 3.3 en Google Cloud Run"),
    (TEAL,  "Base de Datos",   "PostgreSQL 15 en Supabase (cloud)"),
    (PURP,  "App Movil",       "React Native, Android & iOS"),
]
for i, (col, title, detail) in enumerate(comps):
    ty = Inches(1.85 + i * 1.25)
    R(sl, col, Inches(0.55), ty, Inches(0.05), Inches(0.9))
    T(sl, title, Inches(0.75), ty + Inches(0.02), Inches(4.7), Inches(0.45),
      sz=17, bold=True, color=WHITE)
    T(sl, detail, Inches(0.75), ty + Inches(0.45), Inches(4.7), Inches(0.42),
      sz=13, color=LGRAY)
    if i < len(comps) - 1:
        R(sl, DGRAY, Inches(0.55), ty + Inches(1.0), Inches(4.8), Inches(0.01))

T(sl, "Toda la comunicacion usa HTTPS + JWT",
  Inches(0.55), Inches(6.95), Inches(5.0), Inches(0.38),
  sz=12, italic=True, color=GOLD)
add_transition(sl, "push")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BASE DE DATOS
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 4)
header_block(sl, "Datos", "Base de Datos",
             "Supabase - PostgreSQL 15 - Esquema relacional")

# Numero hero
T(sl, "15", Inches(0.55), Inches(1.9), Inches(2.5), Inches(2.5),
  sz=110, bold=True, color=GOLD)
gold_line(sl, Inches(0.55), Inches(4.25), Inches(2.3), Inches(0.04))
T(sl, "TABLAS", Inches(0.55), Inches(4.35), Inches(2.3), Inches(0.5),
  sz=15, bold=True, color=LGRAY)

# Lista de tablas
table_list = [
    "clubs  -  configuracion  -  usuarios",
    "equipos  -  equipos_entrenadores",
    "jugadores  -  permisos_usuario",
    "partidos  -  convocatorias",
    "alineaciones  -  eventos_partido",
    "entrenamientos  -  asistencia_entrenamientos",
    "historial_medico  -  objetivos",
]
tb = sl.shapes.add_textbox(Inches(3.0), Inches(2.0), Inches(3.5), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
p0 = tf.paragraphs[0]; r0 = p0.add_run()
r0.text = "Tablas del modelo:"; r0.font.name = F
r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = GRAY
for t_name in table_list:
    p2 = tf.add_paragraph(); p2.space_before = Pt(9)
    r2 = p2.add_run(); r2.text = " " + t_name
    r2.font.name = F; r2.font.size = Pt(13); r2.font.color.rgb = LGRAY

# Diagrama E/R (carta blanca)
R(sl, WCARD, Inches(6.7), Inches(1.8), Inches(6.2), Inches(5.35))
if not IMG(sl, ER, Inches(6.75), Inches(1.85), Inches(6.1), Inches(5.25)):
    T(sl, "[ E_R.drawio.png ]",
      Inches(6.7), Inches(4.2), Inches(6.2), Inches(0.5),
      sz=11, color=GRAY, align=PP_ALIGN.CENTER)
add_transition(sl, "push")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECNOLOGIAS
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 5)
header_block(sl, "Stack", "Tecnologias Utilizadas")

layers = [
    {
        "col": BLUE, "tag": "CLIENTE ESCRITORIO",
        "items": [
            ("Java 24 + JavaFX 25",    "Lenguaje y framework de interfaz"),
            ("FXML + CSS",             "Vistas declarativas, estilos separados"),
            ("HttpClient Java 11",     "Comunicacion HTTP sin dependencias"),
            ("Gson 2.10",             "Serializacion / deserializacion JSON"),
            ("JasperReports 7.0",      "Generacion de informes en PDF"),
            ("jpackage",               "Instalador nativo ScouteoSetup.exe"),
        ]
    },
    {
        "col": DGRN, "tag": "BACKEND / API",
        "items": [
            ("Spring Boot 3.3",        "Framework principal del backend"),
            ("Spring Security + JWT",  "Autenticacion y autorizacion"),
            ("Spring Data JPA",        "ORM, repositorios, Hibernate"),
            ("BCrypt",                 "Cifrado de contrasenas"),
            ("Supabase PostgreSQL 15", "Base de datos en la nube"),
            ("Google Cloud Run",       "Despliegue serverless del backend"),
        ]
    },
    {
        "col": PURP, "tag": "APP MOVIL",
        "items": [
            ("React Native",           "Framework cross-platform"),
            ("Android & iOS",          "Compilado para ambas plataformas"),
            ("Misma API REST",         "Reutiliza el backend existente"),
            ("JWT",                    "Misma capa de seguridad"),
        ]
    },
]

col_w = Inches(4.0)
for ci, layer in enumerate(layers):
    lx = Inches(0.35 + ci * 4.35)
    ty_header = Inches(1.9)
    # Header de columna
    R(sl, layer["col"], lx, ty_header, col_w, Inches(0.5))
    T(sl, layer["tag"], lx, ty_header + Inches(0.06), col_w, Inches(0.4),
      sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Items
    for ri, (tech, desc) in enumerate(layer["items"]):
        ity = ty_header + Inches(0.6) + ri * Inches(0.8)
        # Separador
        if ri > 0:
            R(sl, DGRAY, lx, ity - Inches(0.08), col_w, Inches(0.01))
        T(sl, tech, lx + Inches(0.15), ity, col_w - Inches(0.2), Inches(0.38),
          sz=13, bold=True, color=WHITE)
        T(sl, desc, lx + Inches(0.15), ity + Inches(0.36), col_w - Inches(0.2), Inches(0.35),
          sz=11, color=GRAY)
add_transition(sl, "push")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DEMO EN VIVO
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 6)
header_block(sl, "Demo", "Flujo End-to-End en Vivo",
             "Se muestra el recorrido completo: BD -> API -> Escritorio -> Movil")

# 5 pasos como circulos conectados
steps = [
    (BLUE,  "1", "Login",         "Dashboard\nJugadores"),
    (GREEN, "2", "Editar",        "Cambiar dato\nGuardar"),
    (LBLUE, "3", "API Cloud",     "PUT /api\nSpring Boot"),
    (GOLD,  "4", "Supabase",      "Tabla BD\nen tiempo real"),
    (PURP,  "5", "App Movil",     "Mismo dato\nreflejado"),
]

# Calcular posiciones
n = len(steps)
cir_d = Inches(1.05)
gap   = Inches(0.45)
total = n * cir_d + (n - 1) * gap
start = (W - total) / 2
cy    = Inches(3.7)

for i, (col, num, title, detail) in enumerate(steps):
    cx = start + i * (cir_d + gap) + cir_d / 2
    # Linea conectora (antes del circulo)
    if i > 0:
        prev_cx = start + (i-1) * (cir_d + gap) + cir_d
        R(sl, DGRAY, prev_cx, cy + cir_d/2 - Inches(0.02),
          cx - cir_d/2 - prev_cx, Inches(0.03))
    # Circulo
    O(sl, col, cx - cir_d/2, cy, cir_d, cir_d)
    T(sl, num, cx - cir_d/2, cy + Inches(0.12), cir_d, cir_d - Inches(0.15),
      sz=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Titulo bajo circulo
    T(sl, title, cx - Inches(0.9), cy + cir_d + Inches(0.15), Inches(1.8), Inches(0.45),
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Detalle
    T(sl, detail, cx - Inches(0.9), cy + cir_d + Inches(0.58), Inches(1.8), Inches(0.75),
      sz=11, color=LGRAY, align=PP_ALIGN.CENTER)

# Nota
R(sl, BG2, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.5), lc=DGRAY, lw=Pt(0.5))
T(sl, "Esta diapositiva permanece en pantalla mientras se realiza la demostracion en directo",
  Inches(0.55), Inches(6.72), Inches(12.23), Inches(0.45),
  sz=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_transition(sl, "cover")

# Que mas mostrar (row de pills)
also = ["Formulario edicion jugador", "Partidos + Convocatoria", "Estadisticas jugador", "Generacion PDF"]
pill_w = Inches(2.9)
pill_start = (W - len(also) * pill_w - (len(also)-1) * Inches(0.1)) / 2
for i, label in enumerate(also):
    lx2 = pill_start + i * (pill_w + Inches(0.1))
    R(sl, DGRAY, lx2, Inches(6.05), pill_w, Inches(0.46))
    T(sl, label, lx2, Inches(6.07), pill_w, Inches(0.42),
      sz=12, color=LGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl)
branding(sl); slide_num(sl, 7)
header_block(sl, "Resultados", "Conclusiones")

# --- Columna izquierda: Propuesto ---
R(sl, BG2, Inches(0.45), Inches(1.95), Inches(5.85), Inches(5.2), lc=DGRAY, lw=Pt(0.5))
R(sl, BLUE, Inches(0.45), Inches(1.95), Inches(5.85), Inches(0.05))
T(sl, "Lo que propuse", Inches(0.45), Inches(2.05), Inches(5.85), Inches(0.5),
  sz=14, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)
R(sl, DGRAY, Inches(0.45), Inches(2.52), Inches(5.85), Inches(0.01))

propuesto = [
    "App escritorio JavaFX multiplataforma",
    "API REST en la nube",
    "Base de datos relacional centralizada",
    "Autenticacion JWT + control de roles",
    "Generacion de informes PDF",
    "App movil complementaria",
]
tb = sl.shapes.add_textbox(Inches(0.65), Inches(2.6), Inches(5.5), Inches(4.4))
tf = tb.text_frame; tf.word_wrap = True
p0 = tf.paragraphs[0]; r0 = p0.add_run(); r0.text = ""
for item in propuesto:
    p = tf.add_paragraph(); p.space_before = Pt(14)
    r = p.add_run(); r.text = "  " + item
    r.font.name = F; r.font.size = Pt(14); r.font.color.rgb = LGRAY

# --- Columna derecha: Entregado ---
R(sl, BG2, Inches(6.75), Inches(1.95), Inches(6.13), Inches(5.2), lc=DGRAY, lw=Pt(0.5))
R(sl, GREEN, Inches(6.75), Inches(1.95), Inches(6.13), Inches(0.05))
T(sl, "Lo que entrego", Inches(6.75), Inches(2.05), Inches(6.13), Inches(0.5),
  sz=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
R(sl, DGRAY, Inches(6.75), Inches(2.52), Inches(6.13), Inches(0.01))

entregado = [
    (GREEN, True,  "Instalador .exe  -  JavaFX 25  -  Windows"),
    (GREEN, True,  "Spring Boot en Google Cloud Run"),
    (GREEN, True,  "15 tablas en Supabase PostgreSQL"),
    (GREEN, True,  "3 roles  -  Spring Security  -  BCrypt"),
    (GREEN, True,  "JasperReports  -  3 tipos de informe"),
    (GOLD,  False, "React Native  -  funcionalidades basicas"),
]
tb2 = sl.shapes.add_textbox(Inches(6.95), Inches(2.6), Inches(5.8), Inches(4.4))
tf2 = tb2.text_frame; tf2.word_wrap = True
p0r = tf2.paragraphs[0]; r0r = p0r.add_run(); r0r.text = ""
for col, ok, item in entregado:
    p = tf2.add_paragraph(); p.space_before = Pt(12)
    r = p.add_run()
    r.text = ("OK  " if ok else "~   ") + item
    r.font.name = F; r.font.size = Pt(14); r.font.color.rgb = col

# Divisor central
R(sl, GOLD, Inches(6.5), Inches(1.95), Inches(0.05), Inches(5.2))

# Mejoras futuras
R(sl, BG2, Inches(0.45), Inches(7.1), Inches(12.43), Inches(0.32), lc=DGRAY, lw=Pt(0.5))
T(sl, "Mejoras futuras:  modo offline  -  notificaciones push  -  modulo economico  -  IA de alineaciones",
  Inches(0.45), Inches(7.12), Inches(12.43), Inches(0.3),
  sz=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_transition(sl, "fade")


# ── Guardar ─────────────────────────────────────────────────────────────────
out = r"C:\Users\User\Desktop\Scouteo_inicios\SCOUTEO_TFG.pptx"
prs.save(out)
print(f"OK -> {out}")
